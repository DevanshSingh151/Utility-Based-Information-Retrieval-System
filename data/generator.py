import json
import os
import random
from datetime import datetime, timedelta
import PyPDF2
from pptx import Presentation
import textwrap

def split_into_chunks(text, words_per_chunk=300):
    words = text.split()
    chunks = []
    for i in range(0, len(words), words_per_chunk):
        chunk = " ".join(words[i:i + words_per_chunk])
        if len(chunk.split()) > 50:  # Minimum 50 words to be a valid chunk
            chunks.append(chunk)
    return chunks

def extract_text_from_pdf(pdf_path):
    text = ""
    try:
        with open(pdf_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                extracted = page.extract_text()
                if extracted:
                    text += extracted + " "
    except Exception as e:
        print(f"Error reading PDF {pdf_path}: {e}")
    return text

def extract_text_from_pptx(pptx_path):
    text = ""
    try:
        prs = Presentation(pptx_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + " "
    except Exception as e:
        print(f"Error reading PPTX {pptx_path}: {e}")
    return text

def determine_domain(filename):
    fname = filename.lower()
    if 'os' in fname:
        return 'operating_systems', 'OS'
    elif 'db' in fname:
        return 'dbms', 'DBMS'
    elif 'cao' in fname:
        return 'computer_architecture', 'CAO'
    elif 'bsts' in fname or 'sts' in fname:
        return 'soft_skills', 'STS'
    else:
        # Fallback 
        return 'general', 'GEN'

def generate_sample_data():
    base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    docs_dir = os.path.join(base_dir, 'documents')
    data_dir = os.path.join(base_dir, 'data')
    output_path = os.path.join(data_dir, 'documents.json')
    
    if not os.path.exists(data_dir):
        os.makedirs(data_dir)
        
    print(f"Reading local documents from {docs_dir}...")
    documents = []
    
    domain_counts = {}
    
    if os.path.exists(docs_dir):
        files = os.listdir(docs_dir)
        for f in files:
            file_path = os.path.join(docs_dir, f)
            domain, prefix = determine_domain(f)
            
            if domain not in domain_counts:
                domain_counts[domain] = 0
            
            text = ""
            if f.endswith('.pdf'):
                text = extract_text_from_pdf(file_path)
            elif f.endswith('.pptx'):
                text = extract_text_from_pptx(file_path)
                
            if text:
                chunks = split_into_chunks(text, words_per_chunk=300)
                for i, chunk in enumerate(chunks):
                    doc_id = f"{domain}_{domain_counts[domain] + 1:03d}"
                    doc = {
                        "id": doc_id,
                        "title": f"Excerpt from {f} (Part {i+1})",
                        "content": chunk,
                        "source": "academic",
                        "source_file": f,
                        "date": (datetime.now() - timedelta(days=random.randint(1, 1500))).strftime('%Y-%m-%d'),
                        "readability_score": round(random.uniform(40.0, 70.0), 1),
                        "domain": domain
                    }
                    documents.append(doc)
                    domain_counts[domain] += 1
    
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(documents, f, indent=4)
        
    print(f"Successfully generated {len(documents)} chunks from real documents and saved to {output_path}")
    for domain, count in domain_counts.items():
        print(f" - {domain}: {count} chunks")

if __name__ == "__main__":
    generate_sample_data()


if __name__ == "__main__":
    generate_sample_data()
